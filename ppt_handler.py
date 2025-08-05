"""
Optimized PowerPoint document handling and text frame updates
"""
import logging
import re
from typing import List, Dict, Any, Tuple, Optional
from dataclasses import dataclass
from pathlib import Path
from pptx.dml.color import RGBColor
from config import Config
from dependencies import DependencyManager
from translation_engine import TranslationEngine
from text_utils import SlideTextCollector

logger = logging.getLogger(__name__)


@dataclass
class TranslationResult:
    """Data class for translation results"""
    translated_count: int = 0
    translated_notes_count: int = 0
    total_shapes: int = 0
    errors: List[str] = None
    
    def __post_init__(self):
        if self.errors is None:
            self.errors = []


class FormattingExtractor:
    """Handles extraction of formatting information from PowerPoint elements"""
    
    @staticmethod
    def extract_paragraph_structure(text_frame):
        """Extract paragraph structure including bullets, indentation, margins, and formatting"""
        paragraph_info = []
        
        try:
            for paragraph in text_frame.paragraphs:
                info = FormattingExtractor._extract_single_paragraph_info(paragraph)
                paragraph_info.append(info)
        except Exception as e:
            logger.error(f"Error extracting paragraph structure: {e}")
        
        return paragraph_info
    
    @staticmethod
    def _extract_single_paragraph_info(paragraph):
        """Extract information from a single paragraph"""
        info = {
            'level': getattr(paragraph, 'level', 0),
            'text': paragraph.text,
            'runs': [],
            'bullet_format': None,
            'alignment': getattr(paragraph, 'alignment', None),
            'space_before': getattr(paragraph, 'space_before', None),
            'space_after': getattr(paragraph, 'space_after', None),
            'line_spacing': getattr(paragraph, 'line_spacing', None),
            'margin_left': None,
            'indent': None
        }
        
        # Extract XML-based formatting
        FormattingExtractor._extract_xml_formatting(paragraph, info)
        
        # Extract run formatting
        for run in paragraph.runs:
            if run.text.strip():
                run_info = FormattingExtractor._extract_run_info(run)
                info['runs'].append(run_info)
        
        return info
    
    @staticmethod
    def _extract_xml_formatting(paragraph, info):
        """Extract formatting from paragraph XML"""
        try:
            if not (hasattr(paragraph, '_element') and paragraph._element is not None):
                return
                
            pPr = paragraph._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            if pPr is None:
                return
            
            # Extract margin and indent
            for attr, key in [('marL', 'margin_left'), ('indent', 'indent'), ('algn', 'alignment_xml')]:
                value = pPr.get(attr)
                if value:
                    info[key] = value
            
            # Extract bullet format
            info['bullet_format'] = FormattingExtractor._extract_bullet_format(pPr, paragraph)
            
        except Exception as e:
            logger.debug(f"Could not extract paragraph XML info: {e}")
            # Fallback for indented paragraphs
            if paragraph.level > 0:
                info['bullet_format'] = {'type': 'char', 'char': '•'}
    
    @staticmethod
    def _extract_bullet_format(pPr, paragraph):
        """Extract bullet format from paragraph properties"""
        bullet_elements = [
            ('buNone', 'none'),
            ('buChar', 'char'),
            ('buAutoNum', 'autonum')
        ]
        
        for elem_name, bullet_type in bullet_elements:
            elem = pPr.find(f'.//{{{FormattingExtractor._get_namespace()}}}{elem_name}')
            if elem is not None:
                return FormattingExtractor._create_bullet_format(elem, bullet_type)
        
        # Default for indented paragraphs
        if paragraph.level > 0:
            return {'type': 'char', 'char': '•'}
        
        return None
    
    @staticmethod
    def _create_bullet_format(elem, bullet_type):
        """Create bullet format dictionary from XML element"""
        if bullet_type == 'none':
            return {'type': 'none'}
        elif bullet_type == 'char':
            return {'type': 'char', 'char': elem.get('char', '•')}
        elif bullet_type == 'autonum':
            return {
                'type': 'autonum',
                'num_type': elem.get('type', 'arabicPeriod'),
                'start_at': elem.get('startAt', '1')
            }
        return None
    
    @staticmethod
    def _extract_run_info(run):
        """Extract information from a single run"""
        run_info = {
            'text': run.text,
            'formatting': FormattingExtractor._extract_run_formatting(run)
        }
        
        # Check for hyperlinks
        try:
            if hasattr(run, 'hyperlink') and run.hyperlink:
                if hasattr(run.hyperlink, 'address') and run.hyperlink.address:
                    run_info['hyperlink'] = run.hyperlink.address
        except Exception:
            pass
        
        return run_info
    
    @staticmethod
    def _extract_run_formatting(run) -> Dict[str, Any]:
        """Extract formatting from a run safely"""
        formatting = {
            'font_name': None,
            'font_size': None,
            'font_bold': None,
            'font_italic': None,
            'font_color': None
        }
        
        try:
            if not (hasattr(run, 'font') and run.font):
                return formatting
                
            font = run.font
            
            # Extract basic font properties
            for attr, key in [('name', 'font_name'), ('size', 'font_size'), 
                             ('bold', 'font_bold'), ('italic', 'font_italic')]:
                if hasattr(font, attr):
                    value = getattr(font, attr)
                    if value is not None:
                        formatting[key] = value
            
            # Extract color
            formatting['font_color'] = FormattingExtractor._extract_font_color(font)
            
        except Exception as e:
            logger.debug(f"Could not extract run formatting: {e}")
        
        return formatting
    
    @staticmethod
    def _extract_font_color(font):
        """Extract font color information"""
        try:
            if not (hasattr(font, 'color') and font.color):
                return None
                
            color_obj = font.color
            if hasattr(color_obj, 'type'):
                if color_obj.type == 1 and hasattr(color_obj, 'rgb'):  # RGB
                    return ('rgb', str(color_obj.rgb))
                elif color_obj.type == 2 and hasattr(color_obj, 'theme_color'):  # Theme
                    return ('theme', color_obj.theme_color)
        except Exception:
            pass
        return None
    
    @staticmethod
    def _get_namespace():
        """Get the OpenXML namespace for drawing ML"""
        return "http://schemas.openxmlformats.org/drawingml/2006/main"


class FormattingApplier:
    """Handles application of formatting to PowerPoint elements"""
    
    @staticmethod
    def apply_paragraph_structure(paragraph, para_info, new_text: str):
        """Apply paragraph structure and formatting"""
        try:
            # Clear paragraph content
            paragraph.clear()
            
            # Apply paragraph-level properties
            FormattingApplier._apply_paragraph_properties(paragraph, para_info)
            
            # Apply text with run formatting
            if para_info and para_info.get('runs'):
                FormattingApplier._apply_runs_with_formatting(paragraph, new_text, para_info['runs'])
            else:
                # Simple text
                run = paragraph.add_run()
                run.text = new_text
                if para_info and para_info.get('runs'):
                    FormattingApplier._apply_run_formatting(run, para_info['runs'][0]['formatting'])
                    
        except Exception as e:
            logger.error(f"Failed to apply paragraph structure: {e}")
            # Fallback
            paragraph.clear()
            paragraph.add_run().text = new_text
    
    @staticmethod
    def _apply_paragraph_properties(paragraph, para_info):
        """Apply paragraph-level properties"""
        if not para_info:
            return
            
        # Set level (indentation)
        if para_info.get('level') is not None:
            paragraph.level = para_info['level']
        
        # Apply XML-based formatting
        FormattingApplier._apply_xml_formatting(paragraph, para_info)
        
        # Apply other properties
        for prop, key in [('alignment', 'alignment'), ('space_before', 'space_before'),
                         ('space_after', 'space_after'), ('line_spacing', 'line_spacing')]:
            if para_info.get(key) is not None:
                try:
                    setattr(paragraph, prop, para_info[key])
                except Exception:
                    pass
        
        # Apply bullet formatting
        FormattingApplier._apply_bullet_format(paragraph, para_info.get('bullet_format'))
    
    @staticmethod
    def _apply_xml_formatting(paragraph, para_info):
        """Apply XML-based formatting (margins, indents)"""
        try:
            if not (hasattr(paragraph, '_element') and paragraph._element is not None):
                return
                
            pPr = paragraph._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            if pPr is None:
                from lxml import etree
                pPr = etree.SubElement(paragraph._element, 
                                     '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            
            # Apply margin and indent
            for xml_attr, info_key in [('marL', 'margin_left'), ('indent', 'indent'), ('algn', 'alignment_xml')]:
                if para_info.get(info_key):
                    pPr.set(xml_attr, para_info[info_key])
                    
        except Exception as e:
            logger.debug(f"Could not apply XML formatting: {e}")
    
    @staticmethod
    def _apply_bullet_format(paragraph, bullet_format):
        """Apply bullet formatting to paragraph"""
        if not bullet_format:
            return
        
        try:
            FormattingApplier._apply_bullet_xml(paragraph, bullet_format)
        except Exception as e:
            logger.debug(f"Could not apply bullet format: {e}")
    
    @staticmethod
    def _apply_bullet_xml(paragraph, bullet_format):
        """Apply bullet formatting via XML manipulation"""
        if not (hasattr(paragraph, '_element') and paragraph._element is not None):
            return
            
        pPr = paragraph._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(paragraph._element, 
                                 '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        
        # Remove existing bullet elements
        namespace = FormattingExtractor._get_namespace()
        for elem_name in ['buNone', 'buChar', 'buAutoNum']:
            for elem in pPr.findall(f'.//{{{namespace}}}{elem_name}'):
                pPr.remove(elem)
        
        # Add new bullet element
        bullet_type = bullet_format.get('type')
        if bullet_type == 'none':
            FormattingApplier._add_bullet_none(pPr, namespace)
        elif bullet_type == 'char':
            FormattingApplier._add_bullet_char(pPr, namespace, bullet_format.get('char', '•'))
        elif bullet_type == 'autonum':
            FormattingApplier._add_bullet_autonum(pPr, namespace, bullet_format)
    
    @staticmethod
    def _add_bullet_none(pPr, namespace):
        """Add buNone element"""
        from lxml import etree
        etree.SubElement(pPr, f'{{{namespace}}}buNone')
    
    @staticmethod
    def _add_bullet_char(pPr, namespace, char):
        """Add buChar element"""
        from lxml import etree
        buChar = etree.SubElement(pPr, f'{{{namespace}}}buChar')
        buChar.set('char', char)
    
    @staticmethod
    def _add_bullet_autonum(pPr, namespace, bullet_format):
        """Add buAutoNum element"""
        from lxml import etree
        buAutoNum = etree.SubElement(pPr, f'{{{namespace}}}buAutoNum')
        buAutoNum.set('type', bullet_format.get('num_type', 'arabicPeriod'))
        start_at = bullet_format.get('start_at', '1')
        if start_at != '1':
            buAutoNum.set('startAt', start_at)
    
    @staticmethod
    def _apply_runs_with_formatting(paragraph, new_text: str, run_info_list):
        """Apply text with preserved run formatting"""
        try:
            if len(run_info_list) == 1:
                # Single run - simple case
                run = paragraph.add_run()
                run.text = new_text
                FormattingApplier._apply_run_formatting(run, run_info_list[0]['formatting'])
            else:
                # Multiple runs - preserve special formatting
                FormattingApplier._apply_multiple_runs(paragraph, new_text, run_info_list)
                
        except Exception as e:
            logger.error(f"Error applying runs with formatting: {e}")
            # Fallback
            if not paragraph.runs:
                run = paragraph.add_run()
                run.text = new_text
                if run_info_list:
                    FormattingApplier._apply_run_formatting(run, run_info_list[0]['formatting'])
    
    @staticmethod
    def _apply_multiple_runs(paragraph, new_text: str, run_info_list):
        """Apply multiple runs with different formatting"""
        remaining_text = new_text
        
        for run_info in run_info_list:
            original_text = run_info['text'].strip()
            
            if original_text and original_text in remaining_text:
                parts = remaining_text.split(original_text, 1)
                
                # Text before this run
                if parts[0]:
                    run = paragraph.add_run()
                    run.text = parts[0]
                    FormattingApplier._apply_run_formatting(run, run_info_list[0]['formatting'])
                
                # This run with its formatting
                run = paragraph.add_run()
                run.text = original_text
                FormattingApplier._apply_run_formatting(run, run_info['formatting'])
                
                # Apply hyperlink if present
                if 'hyperlink' in run_info:
                    try:
                        run.hyperlink.address = run_info['hyperlink']
                    except Exception:
                        pass
                
                # Update remaining text
                remaining_text = parts[1] if len(parts) > 1 else ""
        
        # Add any remaining text
        if remaining_text:
            run = paragraph.add_run()
            run.text = remaining_text
            FormattingApplier._apply_run_formatting(run, run_info_list[-1]['formatting'])
        
        # If no runs were added, add the whole text
        if not paragraph.runs:
            run = paragraph.add_run()
            run.text = new_text
            FormattingApplier._apply_run_formatting(run, run_info_list[0]['formatting'])
    
    @staticmethod
    def _apply_run_formatting(run, formatting: Dict[str, Any]):
        """Apply formatting to a run safely"""
        try:
            if not (hasattr(run, 'font') and run.font):
                return
                
            font = run.font
            
            # Apply basic properties
            for key, attr in [('font_name', 'name'), ('font_size', 'size'), 
                             ('font_bold', 'bold'), ('font_italic', 'italic')]:
                if formatting.get(key) is not None:
                    setattr(font, attr, formatting[key])
            
            # Apply color
            FormattingApplier._apply_font_color(font, formatting.get('font_color'))
            
        except Exception as e:
            logger.debug(f"Could not apply run formatting: {e}")
    
    @staticmethod
    def _apply_font_color(font, color_info):
        """Apply font color information"""
        if not color_info or not isinstance(color_info, tuple) or len(color_info) != 2:
            return
            
        try:
            color_type, color_value = color_info
            if color_type == 'rgb' and color_value:
                FormattingApplier._apply_rgb_color(font, color_value)
            elif color_type == 'theme':
                font.color.theme_color = color_value
        except Exception:
            pass
    
    @staticmethod
    def _apply_rgb_color(font, color_value):
        """Apply RGB color to font"""
        if isinstance(color_value, str) and len(color_value) == 6:
            rgb_int = int(color_value, 16)
            font.color.rgb = RGBColor(
                (rgb_int >> 16) & 0xFF,
                (rgb_int >> 8) & 0xFF,
                rgb_int & 0xFF
            )


class TextFrameUpdater:
    """Handles updating PowerPoint text frames with translations"""
    
    @staticmethod
    def update_text_frame(text_frame, new_text: str):
        """Update text frame while preserving formatting, bullets, and indentation"""
        try:
            if not text_frame.paragraphs:
                text_frame.text = new_text
                return
            
            # Extract paragraph structure information
            paragraph_info = FormattingExtractor.extract_paragraph_structure(text_frame)
            
            # Check for hyperlinks
            has_hyperlinks = TextFrameUpdater._has_hyperlinks(text_frame)
            
            if has_hyperlinks:
                logger.debug("Hyperlinks detected, using safe hyperlink preservation")
                TextFrameUpdater._update_with_hyperlinks_safe(text_frame, new_text, paragraph_info)
                return
            
            # Choose update strategy based on structure
            TextFrameUpdater._choose_update_strategy(text_frame, new_text, paragraph_info)
                
        except Exception as e:
            logger.error(f"Formatting error: {str(e)}")
            text_frame.text = new_text
    
    @staticmethod
    def _choose_update_strategy(text_frame, new_text: str, paragraph_info):
        """Choose the appropriate update strategy"""
        new_lines = new_text.strip().split('\n')
        
        # Single paragraph case
        if len(text_frame.paragraphs) == 1 and len(new_lines) == 1:
            para_info = paragraph_info[0] if paragraph_info else None
            FormattingApplier.apply_paragraph_structure(text_frame.paragraphs[0], para_info, new_text.strip())
            return
        
        # Multiple paragraphs with same count
        if len(new_lines) == len(text_frame.paragraphs):
            TextFrameUpdater._update_matching_paragraphs(text_frame, new_lines, paragraph_info)
        else:
            # Different structure - rebuild with preserved formatting
            TextFrameUpdater._rebuild_with_structure(text_frame, new_text, paragraph_info)
    
    @staticmethod
    def _update_matching_paragraphs(text_frame, new_lines, paragraph_info):
        """Update paragraphs when counts match"""
        for i, (paragraph, new_line) in enumerate(zip(text_frame.paragraphs, new_lines)):
            if new_line.strip():
                para_info = paragraph_info[i] if i < len(paragraph_info) else None
                FormattingApplier.apply_paragraph_structure(paragraph, para_info, new_line.strip())
    
    @staticmethod
    def _rebuild_with_structure(text_frame, new_text: str, paragraph_info):
        """Rebuild text frame with preserved structure"""
        try:
            text_frame.clear()
            new_lines = new_text.strip().split('\n')
            
            for i, line in enumerate(new_lines):
                if i > 0:
                    paragraph = text_frame.add_paragraph()
                else:
                    paragraph = text_frame.paragraphs[0]
                
                # Use corresponding paragraph info if available
                para_info = paragraph_info[i] if i < len(paragraph_info) else (paragraph_info[0] if paragraph_info else None)
                FormattingApplier.apply_paragraph_structure(paragraph, para_info, line.strip())
                
        except Exception as e:
            logger.error(f"Structure rebuild failed: {e}")
            text_frame.text = new_text
    
    @staticmethod
    def _update_with_hyperlinks_safe(text_frame, new_text: str, paragraph_info=None):
        """Update text frame while preserving hyperlinks and structure"""
        try:
            new_lines = new_text.strip().split('\n')
            
            for i, line in enumerate(new_lines):
                if i < len(text_frame.paragraphs):
                    paragraph = text_frame.paragraphs[i]
                else:
                    paragraph = text_frame.add_paragraph()
                
                # Get paragraph info
                para_info = paragraph_info[i] if paragraph_info and i < len(paragraph_info) else None
                
                # Clear and apply structure
                paragraph.clear()
                
                if para_info:
                    FormattingApplier._apply_paragraph_properties(paragraph, para_info)
                    TextFrameUpdater._apply_hyperlinks_to_paragraph(paragraph, line.strip(), para_info)
                else:
                    paragraph.add_run().text = line.strip()
                    
        except Exception as e:
            logger.error(f"Safe hyperlink preservation failed: {e}")
            text_frame.text = new_text
    
    @staticmethod
    def _apply_hyperlinks_to_paragraph(paragraph, line: str, para_info):
        """Apply hyperlinks to paragraph with structure preservation"""
        try:
            runs_info = para_info.get('runs', [])
            hyperlink_runs = [run for run in runs_info if run.get('hyperlink')]
            
            if not hyperlink_runs:
                # No hyperlinks, just add text with formatting
                run = paragraph.add_run()
                run.text = line
                if runs_info:
                    FormattingApplier._apply_run_formatting(run, runs_info[0]['formatting'])
                return
            
            # Apply hyperlinks
            remaining_text = line
            
            for hyperlink_run in hyperlink_runs:
                original_text = hyperlink_run['text'].strip()
                hyperlink_url = hyperlink_run['hyperlink']
                
                # Find hyperlink text in translated line
                hyperlink_text = TextFrameUpdater._find_hyperlink_text(remaining_text, original_text)
                
                if hyperlink_text and hyperlink_text in remaining_text:
                    parts = remaining_text.split(hyperlink_text, 1)
                    
                    # Text before hyperlink
                    if parts[0]:
                        run = paragraph.add_run()
                        run.text = parts[0]
                        default_formatting = next((r['formatting'] for r in runs_info if not r.get('hyperlink')), 
                                                runs_info[0]['formatting'] if runs_info else {})
                        FormattingApplier._apply_run_formatting(run, default_formatting)
                    
                    # Hyperlink text
                    run = paragraph.add_run()
                    run.text = hyperlink_text
                    FormattingApplier._apply_run_formatting(run, hyperlink_run['formatting'])
                    
                    # Apply hyperlink
                    try:
                        run.hyperlink.address = hyperlink_url
                        logger.debug(f"Applied hyperlink: '{hyperlink_text}' -> {hyperlink_url}")
                    except Exception as e:
                        logger.debug(f"Could not apply hyperlink: {e}")
                    
                    remaining_text = parts[1] if len(parts) > 1 else ""
                    break
            
            # Add remaining text
            if remaining_text:
                run = paragraph.add_run()
                run.text = remaining_text
                default_formatting = next((r['formatting'] for r in runs_info if not r.get('hyperlink')), 
                                        runs_info[0]['formatting'] if runs_info else {})
                FormattingApplier._apply_run_formatting(run, default_formatting)
                
        except Exception as e:
            logger.error(f"Error applying hyperlinks: {e}")
            if not paragraph.runs:
                run = paragraph.add_run()
                run.text = line
                if para_info.get('runs'):
                    FormattingApplier._apply_run_formatting(run, para_info['runs'][0]['formatting'])
    
    @staticmethod
    def _has_hyperlinks(text_frame):
        """Check if text frame contains hyperlinks"""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if hasattr(run, 'hyperlink') and run.hyperlink:
                        if hasattr(run.hyperlink, 'address') and run.hyperlink.address:
                            return True
        except Exception:
            pass
        return False
    
    @staticmethod
    def _find_hyperlink_text(translated_text: str, original_text: str):
        """Find text that should be hyperlinked"""
        # First try exact match
        if original_text in translated_text:
            return original_text
        
        # Common hyperlink patterns
        patterns = [
            'Boto3', 'Code samples', 'Starter Toolkit', 'samples', 'toolkit',
            '코드 샘플', '샘플', '툴킷', '스타터', 'Boto3', '코드'
        ]
        
        words = translated_text.split()
        for pattern in patterns:
            for word in words:
                if pattern.lower() in word.lower() or word.lower() in pattern.lower():
                    return word
        
        # Return first meaningful word
        meaningful_words = [word for word in words if len(word) > 2]
        return meaningful_words[0] if meaningful_words else None


class ComplexityAnalyzer:
    """Analyzes slide complexity to determine translation strategy"""
    
    @staticmethod
    def slide_has_complex_formatting(text_items: List[Dict]) -> bool:
        """Check if slide has complex formatting including bullets and indentation"""
        for item in text_items:
            if item['type'] == 'text_frame_unified':
                if ComplexityAnalyzer._text_frame_has_complex_formatting(item['text_frame']):
                    return True
        return False
    
    @staticmethod
    def _text_frame_has_complex_formatting(text_frame) -> bool:
        """Check if text frame has complex formatting"""
        for paragraph in text_frame.paragraphs:
            # Check for indentation (lists)
            if hasattr(paragraph, 'level') and paragraph.level and paragraph.level > 0:
                logger.debug(f"Found indented paragraph with level: {paragraph.level}")
                return True
            
            # Check for bullet formatting in XML
            if ComplexityAnalyzer._has_bullet_formatting(paragraph):
                logger.debug("Found bullet formatting in XML")
                return True
            
            # Check for multiple runs with different formatting
            if ComplexityAnalyzer._has_multiple_formatting_styles(paragraph):
                logger.debug("Found multiple formatting styles")
                return True
        
        return False
    
    @staticmethod
    def _has_bullet_formatting(paragraph) -> bool:
        """Check if paragraph has bullet formatting"""
        try:
            if not (hasattr(paragraph, '_element') and paragraph._element):
                return False
                
            pPr = paragraph._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            if pPr is None:
                return False
            
            # Check for any bullet formatting elements
            bullet_elements = ['buFont', 'buChar', 'buAutoNum']
            namespace = FormattingExtractor._get_namespace()
            
            for elem_name in bullet_elements:
                if pPr.find(f'.//{{{namespace}}}{elem_name}') is not None:
                    return True
                    
        except Exception as e:
            logger.debug(f"Could not check bullet formatting: {e}")
        
        return False
    
    @staticmethod
    def _has_multiple_formatting_styles(paragraph) -> bool:
        """Check if paragraph has multiple formatting styles"""
        if len(paragraph.runs) <= 1:
            return False
        
        colors = []
        italic_states = []
        
        for run in paragraph.runs:
            try:
                # Check colors
                if hasattr(run.font, 'color') and run.font.color:
                    color = run.font.color
                    if hasattr(color, 'type') and color.type == 1:  # RGB
                        colors.append(str(color.rgb))
                    elif hasattr(color, 'type') and color.type == 2:  # Theme
                        colors.append(f"theme_{color.theme_color}")
                
                # Check italic
                italic_states.append(run.font.italic if hasattr(run.font, 'italic') else None)
                
            except Exception:
                pass
        
        # If we have different colors or italic states
        return len(set(colors)) > 1 or len(set(italic_states)) > 1


class TranslationStrategy:
    """Handles different translation strategies"""
    
    def __init__(self, engine: TranslationEngine, text_updater: TextFrameUpdater):
        self.engine = engine
        self.text_updater = text_updater
    
    def translate_slide(self, slide, target_language: str) -> Tuple[int, bool]:
        """Translate a single slide using appropriate strategy"""
        text_items, notes_text = SlideTextCollector().collect_slide_texts(slide)
        
        translated_count = 0
        notes_translated = False
        
        # Translate notes if present
        if notes_text:
            notes_translated = self._translate_notes(slide, notes_text, target_language)
        
        # Choose translation strategy
        if ComplexityAnalyzer.slide_has_complex_formatting(text_items):
            logger.info("🎨 Complex formatting detected, using individual translation")
            translated_count = self._translate_individually(text_items, target_language)
        elif len(text_items) > Config.CONTEXT_THRESHOLD:
            translated_count = self._translate_with_context(text_items, target_language)
        else:
            translated_count = self._translate_with_batch(text_items, target_language)
        
        return translated_count, notes_translated
    
    def _translate_notes(self, slide, notes_text: str, target_language: str) -> bool:
        """Translate slide notes"""
        try:
            translated_notes = self.engine.translate_text(notes_text, target_language)
            if translated_notes != notes_text:
                slide.notes_slide.notes_text_frame.text = translated_notes
                return True
        except Exception as e:
            logger.error(f"Error translating slide notes: {str(e)}")
        return False
    
    def _translate_individually(self, text_items: List[Dict], target_language: str) -> int:
        """Translate each text individually to preserve formatting"""
        translated_count = 0
        
        logger.info("🎨 Using individual translation to preserve complex formatting")
        
        for i, item in enumerate(text_items):
            try:
                original_text = item['text']
                if not original_text.strip():
                    continue
                    
                logger.debug(f"Translating item {i+1}/{len(text_items)}: '{original_text[:50]}...'")
                translation = self.engine.translate_text(original_text, target_language)
                
                if original_text != translation:
                    if self._apply_translation_to_item(item, translation):
                        translated_count += 1
                        logger.debug(f"✅ Applied translation to {item['type']}")
                        
            except Exception as e:
                logger.error(f"Individual translation failed for item {i}: {str(e)}")
        
        logger.info(f"🎯 Individual translation completed: {translated_count}/{len(text_items)} items translated")
        return translated_count
    
    def _translate_with_context(self, text_items: List[Dict], target_language: str) -> int:
        """Translate using context-aware approach"""
        if not text_items:
            return 0
        
        try:
            translations = self.engine.translate_with_context(text_items, target_language)
            return self._apply_translations(text_items, translations)
        except Exception as e:
            logger.error(f"Context translation failed: {str(e)}")
            return self._translate_with_batch(text_items, target_language)
    
    def _translate_with_batch(self, text_items: List[Dict], target_language: str) -> int:
        """Translate using batch approach"""
        if not text_items:
            return 0
        
        texts_to_translate = [item['text'] for item in text_items]
        translated_count = 0
        
        # Process in batches
        for i in range(0, len(texts_to_translate), Config.BATCH_SIZE):
            batch_items = text_items[i:i + Config.BATCH_SIZE]
            batch_texts = texts_to_translate[i:i + Config.BATCH_SIZE]
            
            try:
                batch_translations = self.engine.translate_batch(batch_texts, target_language)
                translated_count += self._apply_translations(batch_items, batch_translations)
            except Exception as e:
                logger.error(f"Batch translation failed: {str(e)}")
                # Individual fallback
                for item in batch_items:
                    try:
                        translation = self.engine.translate_text(item['text'], target_language)
                        if self._apply_translation_to_item(item, translation):
                            translated_count += 1
                    except Exception:
                        pass
        
        return translated_count
    
    def _apply_translations(self, text_items: List[Dict], translations: List[str]) -> int:
        """Apply translations back to the original shapes"""
        if len(text_items) != len(translations):
            logger.error(f"Translation count mismatch: {len(text_items)} items, {len(translations)} translations")
            return 0
        
        translated_count = 0
        
        for item, translation in zip(text_items, translations):
            if item['text'] != translation:
                if self._apply_translation_to_item(item, translation):
                    translated_count += 1
        
        return translated_count
    
    def _apply_translation_to_item(self, item: Dict, translation: str) -> bool:
        """Apply translation to a single item"""
        try:
            item_type = item['type']
            
            if item_type == 'table_cell':
                cell = item['cell']
                if hasattr(cell, 'text_frame') and cell.text_frame:
                    self.text_updater.update_text_frame(cell.text_frame, translation)
                else:
                    cell.text = translation
                return True
                
            elif item_type == 'text_frame_unified':
                text_frame = item['text_frame']
                self.text_updater.update_text_frame(text_frame, translation)
                return True
                
            elif item_type == 'direct_text':
                item['shape'].text = translation
                return True
            
        except Exception as e:
            logger.error(f"Error applying translation: {str(e)}")
        
        return False


class PowerPointTranslator:
    """Main PowerPoint translation class"""
    
    def __init__(self, model_id: str = Config.DEFAULT_MODEL_ID, enable_polishing: bool = Config.ENABLE_POLISHING):
        self.model_id = model_id
        self.enable_polishing = enable_polishing
        self.engine = TranslationEngine(model_id, enable_polishing)
        self.text_updater = TextFrameUpdater()
        self.strategy = TranslationStrategy(self.engine, self.text_updater)
        self.deps = DependencyManager()
    
    def translate_presentation(self, input_file: str, output_file: str, target_language: str) -> TranslationResult:
        """Translate entire PowerPoint presentation"""
        try:
            Presentation = self.deps.require('pptx')
            prs = Presentation(input_file)
            result = TranslationResult()
            
            total_slides = len(prs.slides)
            logger.info(f"🎯 Starting translation of {total_slides} slides...")
            logger.info(f"🎨 Translation mode: {'Natural/Polished' if self.enable_polishing else 'Literal'}")
            
            for slide_idx, slide in enumerate(prs.slides):
                logger.info(f"📄 Processing slide {slide_idx + 1}/{total_slides}")
                
                translated_count, notes_translated = self.strategy.translate_slide(slide, target_language)
                
                result.translated_count += translated_count
                if notes_translated:
                    result.translated_notes_count += 1
                result.total_shapes += len(slide.shapes)
                
                logger.info(f"✅ Slide {slide_idx + 1}: {translated_count} texts translated")
            
            # Save translated presentation
            prs.save(output_file)
            logger.info(f"🎉 Translation completed: {output_file}")
            logger.info(f"📊 Summary: {result.translated_count} texts, {result.translated_notes_count} notes")
            
            return result
            
        except Exception as e:
            logger.error(f"❌ Translation failed: {str(e)}")
            raise

    def translate_specific_slides(self, input_file: str, output_file: str, target_language: str, slide_numbers: List[int]) -> TranslationResult:
        """Translate specific slides in PowerPoint presentation"""
        try:
            Presentation = self.deps.require('pptx')
            prs = Presentation(input_file)
            result = TranslationResult()
            
            total_slides = len(prs.slides)
            
            # Validate slide numbers
            invalid_slides = [num for num in slide_numbers if num < 1 or num > total_slides]
            if invalid_slides:
                error_msg = f"Invalid slide numbers: {invalid_slides}. Valid range: 1-{total_slides}"
                logger.error(error_msg)
                result.errors.append(error_msg)
                return result
            
            # Remove duplicates and sort
            slide_numbers = sorted(list(set(slide_numbers)))
            
            logger.info(f"🎯 Starting translation of {len(slide_numbers)} specific slides: {slide_numbers}")
            logger.info(f"🎨 Translation mode: {'Natural/Polished' if self.enable_polishing else 'Literal'}")
            
            for slide_num in slide_numbers:
                slide_idx = slide_num - 1  # Convert to 0-based index
                slide = prs.slides[slide_idx]
                
                logger.info(f"📄 Processing slide {slide_num}/{total_slides}")
                
                translated_count, notes_translated = self.strategy.translate_slide(slide, target_language)
                
                result.translated_count += translated_count
                if notes_translated:
                    result.translated_notes_count += 1
                result.total_shapes += len(slide.shapes)
                
                logger.info(f"✅ Slide {slide_num}: {translated_count} texts translated")
            
            # Save translated presentation
            prs.save(output_file)
            logger.info(f"🎉 Translation completed: {output_file}")
            logger.info(f"📊 Summary: {result.translated_count} texts, {result.translated_notes_count} notes from {len(slide_numbers)} slides")
            
            return result
            
        except Exception as e:
            logger.error(f"❌ Translation failed: {str(e)}")
            raise

    def get_slide_count(self, input_file: str) -> int:
        """Get total number of slides in PowerPoint presentation"""
        try:
            Presentation = self.deps.require('pptx')
            prs = Presentation(input_file)
            return len(prs.slides)
        except Exception as e:
            logger.error(f"❌ Failed to get slide count: {str(e)}")
            raise

    def get_slide_preview(self, input_file: str, slide_number: int, max_chars: int = 200) -> str:
        """Get a preview of text content from a specific slide"""
        try:
            Presentation = self.deps.require('pptx')
            prs = Presentation(input_file)
            
            if slide_number < 1 or slide_number > len(prs.slides):
                raise ValueError(f"Invalid slide number: {slide_number}. Valid range: 1-{len(prs.slides)}")
            
            slide = prs.slides[slide_number - 1]  # Convert to 0-based index
            text_items, notes_text = SlideTextCollector().collect_slide_texts(slide)
            
            # Collect all text content
            all_texts = []
            for item in text_items:
                if item['text'].strip():
                    all_texts.append(item['text'].strip())
            
            if notes_text and notes_text.strip():
                all_texts.append(f"[Notes: {notes_text.strip()}]")
            
            # Join and truncate if necessary
            preview = " | ".join(all_texts)
            if len(preview) > max_chars:
                preview = preview[:max_chars] + "..."
            
            return preview if preview else "[No text content found]"
            
        except Exception as e:
            logger.error(f"❌ Failed to get slide preview: {str(e)}")
            raise
